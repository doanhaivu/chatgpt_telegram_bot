import os
import logging
import asyncio
import traceback
import html
import json
import tempfile
import pydub
from pathlib import Path
from datetime import datetime
from datetime import date
import openai
import PyPDF2
import docx
import pptx
import ast
import telegram
from telegram import (
    Update,
    User,
    InlineKeyboardButton,
    InlineKeyboardMarkup,
    BotCommand,
    LabeledPrice, 
    ShippingOption
)
from telegram.ext import (
    Application,
    ApplicationBuilder,
    CallbackContext,
    CommandHandler,
    MessageHandler,
    CallbackQueryHandler,
    AIORateLimiter,
    ContextTypes,
    PreCheckoutQueryHandler,
    ShippingQueryHandler,
    filters
)
from telegram.constants import ParseMode, ChatAction

import config
import database
import openai_utils
import payment

from elasticsearch_utils import ElasticsearchUtils
from utils import RetrievalUtils

# setup
db = database.Database()
logger = logging.getLogger(__name__)
retrieval_utils = RetrievalUtils(config)
es_utils = ElasticsearchUtils(config)

user_semaphores = {}
user_tasks = {}

HELP_MESSAGE = config.help_message

HELP_GROUP_CHAT_MESSAGE = config.help_group_message

PAYMENT_PROVIDER_TOKEN = "284685063:TEST:OTUwYjZiNjJhMjFk"

def split_text_into_chunks(text, chunk_size):
    for i in range(0, len(text), chunk_size):
        yield text[i:i + chunk_size]


async def register_user_if_not_exists(update: Update, context: CallbackContext, user: User):
    if not db.check_if_user_exists(user.id):
        db.add_new_user(
            user.id,
            update.message.chat_id,
            username=user.username,
            first_name=user.first_name,
            last_name= user.last_name
        )
        db.start_new_dialog(user.id)

    if db.get_user_attribute(user.id, "current_dialog_id") is None:
        db.start_new_dialog(user.id)

    if user.id not in user_semaphores:
        user_semaphores[user.id] = asyncio.Semaphore(1)

    if db.get_user_attribute(user.id, "current_model") is None:
        db.set_user_attribute(user.id, "current_model", config.models["available_text_models"][0])

    # back compatibility for n_used_tokens field
    n_used_tokens = db.get_user_attribute(user.id, "n_used_tokens")
    if isinstance(n_used_tokens, int) or isinstance(n_used_tokens, float):  # old format
        new_n_used_tokens = {
            "gpt-3.5-turbo": {
                "n_input_tokens": 0,
                "n_output_tokens": n_used_tokens
            }
        }
        db.set_user_attribute(user.id, "n_used_tokens", new_n_used_tokens)

    # voice message transcription
    if db.get_user_attribute(user.id, "n_transcribed_seconds") is None:
        db.set_user_attribute(user.id, "n_transcribed_seconds", 0.0)

    # image generation
    if db.get_user_attribute(user.id, "n_generated_images") is None:
        db.set_user_attribute(user.id, "n_generated_images", 0)


async def is_bot_mentioned(update: Update, context: CallbackContext):
     try:
         message = update.message

         if message.chat.type == "private":
             return True

         if message.text is not None and ("@" + context.bot.username) in message.text:
             return True

         if message.reply_to_message is not None:
             if message.reply_to_message.from_user.id == context.bot.id:
                 return True
     except:
         return True
     else:
         return False


async def start_handle(update: Update, context: CallbackContext):
    await register_user_if_not_exists(update, context, update.message.from_user)
    user_id = update.message.from_user.id

    db.set_user_attribute(user_id, "last_interaction", datetime.now())
    db.start_new_dialog(user_id)

    reply_text = config.hi_msg
    reply_text += HELP_MESSAGE

    await update.message.reply_text(reply_text, parse_mode=ParseMode.HTML)
    await show_chat_modes_handle(update, context)


async def help_handle(update: Update, context: CallbackContext):
    await register_user_if_not_exists(update, context, update.message.from_user)
    user_id = update.message.from_user.id
    db.set_user_attribute(user_id, "last_interaction", datetime.now())
    await update.message.reply_text(HELP_MESSAGE, parse_mode=ParseMode.HTML)


async def help_group_chat_handle(update: Update, context: CallbackContext):
     await register_user_if_not_exists(update, context, update.message.from_user)
     user_id = update.message.from_user.id
     db.set_user_attribute(user_id, "last_interaction", datetime.now())

     text = HELP_GROUP_CHAT_MESSAGE.format(bot_username="@" + context.bot.username)

     await update.message.reply_text(text, parse_mode=ParseMode.HTML)
     await update.message.reply_video(config.help_group_chat_video_path)


async def retry_handle(update: Update, context: CallbackContext):
    await register_user_if_not_exists(update, context, update.message.from_user)
    if await is_previous_message_not_answered_yet(update, context): return

    user_id = update.message.from_user.id
    db.set_user_attribute(user_id, "last_interaction", datetime.now())

    dialog_messages = db.get_dialog_messages(user_id, dialog_id=None)
    if len(dialog_messages) == 0:
        await update.message.reply_text("No message to retry ü§∑‚Äç‚ôÇÔ∏è")
        return

    last_dialog_message = dialog_messages.pop()
    db.set_dialog_messages(user_id, dialog_messages, dialog_id=None)  # last message was removed from the context

    await message_handle(update, context, message=last_dialog_message["user"], use_new_dialog_timeout=False)


async def message_handle(update: Update, context: CallbackContext, message=None, use_new_dialog_timeout=True):
    # check if bot was mentioned (for group chats)
    if not await is_bot_mentioned(update, context):
        return

    # check if message is edited
    if update.edited_message is not None:
        await edited_message_handle(update, context)
        return

    _message = message or update.message.text

    # remove bot mention (in group chats)
    if update.message.chat.type != "private":
        _message = _message.replace("@" + context.bot.username, "").strip()

    await register_user_if_not_exists(update, context, update.message.from_user)
    if await is_previous_message_not_answered_yet(update, context): return

    user_id = update.message.from_user.id
    chat_mode = db.get_user_attribute(user_id, "current_chat_mode")
    
    if not is_allow_use_api(user_id):
        await update.message.reply_text("You've used up your daily quota. Try again tomorrow or subscribe to our premium plan.")
        return

    if chat_mode == "artist":
        await generate_image_handle(update, context, message=message)
        return

    async def message_handle_fn():
        # new dialog timeout
        if use_new_dialog_timeout:
            if (datetime.now() - db.get_user_attribute(user_id, "last_interaction")).seconds > config.new_dialog_timeout and len(db.get_dialog_messages(user_id)) > 0:
                db.start_new_dialog(user_id)
                await update.message.reply_text(f"Starting new dialog due to timeout (<b>{config.chat_modes[chat_mode]['name']}</b> mode) ‚úÖ", parse_mode=ParseMode.HTML)
        db.set_user_attribute(user_id, "last_interaction", datetime.now())

        
        # in case of CancelledError
        n_input_tokens, n_output_tokens = 0, 0
        current_model = db.get_user_attribute(user_id, "current_model")
        try:
            # send placeholder message to user
            placeholder_message = await update.message.reply_text("...")

            # send typing action
            await update.message.chat.send_action(action="typing")

            if _message is None or len(_message) == 0:
                 await update.message.reply_text("ü•≤ You sent <b>empty message</b>. Please, try again!", parse_mode=ParseMode.HTML)
                 return

            dialog_messages = db.get_dialog_messages(user_id, dialog_id=None)
            parse_mode = {
                "html": ParseMode.HTML,
                "markdown": ParseMode.MARKDOWN
            }[config.chat_modes[chat_mode]["parse_mode"]]

            chatgpt_instance = openai_utils.ChatGPT(model=current_model)
            if config.enable_message_streaming:
                gen = chatgpt_instance.send_message_stream(_message, dialog_messages=dialog_messages, chat_mode=chat_mode)
            else:
                answer, (n_input_tokens, n_output_tokens), n_first_dialog_messages_removed = await chatgpt_instance.send_message(
                    _message,
                    dialog_messages=dialog_messages,
                    chat_mode=chat_mode
                )

                async def fake_gen():
                    yield "finished", answer, (n_input_tokens, n_output_tokens), n_first_dialog_messages_removed

                gen = fake_gen()

            prev_answer = ""
            async for gen_item in gen:
                status, answer, (n_input_tokens, n_output_tokens), n_first_dialog_messages_removed = gen_item

                answer = answer[:4096]  # telegram message limit

                # update only when 100 new symbols are ready
                if abs(len(answer) - len(prev_answer)) < 100 and status != "finished":
                    continue

                try:
                    await context.bot.edit_message_text(answer, chat_id=placeholder_message.chat_id, message_id=placeholder_message.message_id, parse_mode=parse_mode)
                except telegram.error.BadRequest as e:
                    if str(e).startswith("Message is not modified"):
                        continue
                    else:
                        await context.bot.edit_message_text(answer, chat_id=placeholder_message.chat_id, message_id=placeholder_message.message_id)

                await asyncio.sleep(0.01)  # wait a bit to avoid flooding

                prev_answer = answer

            # update user data
            new_dialog_message = {"user": _message, "bot": answer, "date": datetime.now()}
            now = datetime.now()
            current_date = now.strftime("%Y-%m-%d")
            es_utils.put_data("chat-"+current_date, {"user": _message, "bot": answer, "timestamp": now})
            db.set_dialog_messages(
                user_id,
                db.get_dialog_messages(user_id, dialog_id=None) + [new_dialog_message],
                dialog_id=None
            )

            db.update_n_used_tokens(user_id, current_model, n_input_tokens, n_output_tokens)
            db.add_new_user_daily_stats(user_id, str(date.today()), n_input_tokens + n_output_tokens, 1)

        except asyncio.CancelledError:
            # note: intermediate token updates only work when enable_message_streaming=True (config.yml)
            db.update_n_used_tokens(user_id, current_model, n_input_tokens, n_output_tokens)
            db.add_new_user_daily_stats(user_id, str(date.today()), n_input_tokens + n_output_tokens, 1)
            raise

        except Exception as e:
            error_text = f"Something went wrong during completion. Reason: {e}"
            logger.error(error_text)
            await update.message.reply_text(error_text)
            return
        # send message if some messages were removed from the context
        if n_first_dialog_messages_removed > 0:
            if n_first_dialog_messages_removed == 1:
                text = "‚úçÔ∏è <i>Note:</i> Your current dialog is too long, so your <b>first message</b> was removed from the context.\n Send /new command to start new dialog"
            else:
                text = f"‚úçÔ∏è <i>Note:</i> Your current dialog is too long, so <b>{n_first_dialog_messages_removed} first messages</b> were removed from the context.\n Send /new command to start new dialog"
            await update.message.reply_text(text, parse_mode=ParseMode.HTML)

    async with user_semaphores[user_id]:
        task = asyncio.create_task(message_handle_fn())
        user_tasks[user_id] = task

        try:
            await task
        except asyncio.CancelledError:
            await update.message.reply_text("‚úÖ Canceled", parse_mode=ParseMode.HTML)
        else:
            pass
        finally:
            if user_id in user_tasks:
                del user_tasks[user_id]


async def is_previous_message_not_answered_yet(update: Update, context: CallbackContext):
    await register_user_if_not_exists(update, context, update.message.from_user)

    user_id = update.message.from_user.id
    if user_semaphores[user_id].locked():
        text = "‚è≥ Please <b>wait</b> for a reply to the previous message\n"
        text += "Or you can /cancel it"
        await update.message.reply_text(text, reply_to_message_id=update.message.id, parse_mode=ParseMode.HTML)
        return True
    else:
        return False


async def voice_message_handle(update: Update, context: CallbackContext):
    # check if bot was mentioned (for group chats)
    if not await is_bot_mentioned(update, context):
        return

    await register_user_if_not_exists(update, context, update.message.from_user)
    if await is_previous_message_not_answered_yet(update, context): return

    user_id = update.message.from_user.id
    db.set_user_attribute(user_id, "last_interaction", datetime.now())

    voice = update.message.voice
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir = Path(tmp_dir)
        voice_ogg_path = tmp_dir / "voice.ogg"

        # download
        voice_file = await context.bot.get_file(voice.file_id)
        await voice_file.download_to_drive(voice_ogg_path)

        # convert to mp3
        voice_mp3_path = tmp_dir / "voice.mp3"
        pydub.AudioSegment.from_file(voice_ogg_path).export(voice_mp3_path, format="mp3")

        # transcribe
        with open(voice_mp3_path, "rb") as f:
            transcribed_text = await openai_utils.transcribe_audio(f)

            if transcribed_text is None:
                 transcribed_text = ""

    text = f"üé§: <i>{transcribed_text}</i>"
    await update.message.reply_text(text, parse_mode=ParseMode.HTML)

    # update n_transcribed_seconds
    db.set_user_attribute(user_id, "n_transcribed_seconds", voice.duration + db.get_user_attribute(user_id, "n_transcribed_seconds"))

    await message_handle(update, context, message=transcribed_text)

async def image_message_handle(update: Update, context: CallbackContext):
    user_id = await pre_message_handle(update, context)
    return

async def video_message_handle(update: Update, context: CallbackContext):
    user_id = await pre_message_handle(update, context)
    return

async def text_file_message_handle(update: Update, context: CallbackContext):
    user_id = await pre_message_handle(update, context)
    return

async def pre_message_handle(update: Update, context: CallbackContext):
    # check if bot was mentioned (for group chats)
    if not await is_bot_mentioned(update, context):
        return

    await register_user_if_not_exists(update, context, update.message.from_user)
    if await is_previous_message_not_answered_yet(update, context): return

    user_id = update.message.from_user.id
    db.set_user_attribute(user_id, "last_interaction", datetime.now())
    return user_id

async def pdf_message_handle(update: Update, context: CallbackContext):
    user_id = await pre_message_handle(update, context)

    file_id = update.message.document.file_id
    file_name = update.message.document.file_name
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir = Path(tmp_dir)
        pdf_ogg_path = tmp_dir / "pdf.ogg"

        # download
        new_file = await context.bot.get_file(file_id)
        await new_file.download_to_drive(pdf_ogg_path)

        # read
        with open(pdf_ogg_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            meta = reader.metadata
            num_pages = len(reader.pages)
            text = ''
            for i in range(num_pages):
                page_text = reader.pages[i].extract_text()
                text += page_text

    await process_received_content(update, file_name, meta, text, file_id)

async def doc_message_handle(update: Update, context: CallbackContext):
    user_id = await pre_message_handle(update, context)

    file_id = update.message.document.file_id
    file_name = update.message.document.file_name
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir = Path(tmp_dir)
        doc_ogg_path = tmp_dir / "doc.ogg"

        # download
        new_file = await context.bot.get_file(file_id)
        await new_file.download_to_drive(doc_ogg_path)

        # read
        doc = docx.Document(doc_ogg_path)
        meta = doc.core_properties
        text = ''
        for para in doc.paragraphs:
            text += para.text

    await process_received_content(update, file_name, meta, text, file_id)

async def ppt_message_handle(update: Update, context: CallbackContext):
    user_id = await pre_message_handle(update, context)

    file_id = update.message.document.file_id
    file_name = update.message.document.file_name
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir = Path(tmp_dir)
        ppt_ogg_path = tmp_dir / "ppt.ogg"

        # download
        new_file = await context.bot.get_file(file_id)
        await new_file.download_to_drive(ppt_ogg_path)

        # read
        prs = pptx.Presentation(ppt_ogg_path)

        meta = prs.core_properties
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text

    await process_received_content(update, file_name, meta, text, file_id)

async def process_received_content(update, file_name, meta, text, file_id):
    reply = f"Received file \n{file_name}, title: \n{meta.title}, author: \n{meta.author}"
    await update.message.reply_text(reply, parse_mode=ParseMode.HTML)
    response = await retrieval_utils.upsert(text, file_name, meta.title, file_id, meta.author)
    await update.message.reply_text(response, parse_mode=ParseMode.HTML)

async def generate_image_handle(update: Update, context: CallbackContext, message=None):
    await register_user_if_not_exists(update, context, update.message.from_user)
    if await is_previous_message_not_answered_yet(update, context): return

    user_id = update.message.from_user.id
    db.set_user_attribute(user_id, "last_interaction", datetime.now())

    await update.message.chat.send_action(action="upload_photo")

    message = message or update.message.text

    try:
        image_urls = await openai_utils.generate_images(message, n_images=config.return_n_generated_images)
    except openai.error.InvalidRequestError as e:
        if str(e).startswith("Your request was rejected as a result of our safety system"):
            text = "ü•≤ Your request <b>doesn't comply</b> with OpenAI's usage policies.\nWhat did you write there, huh?"
            await update.message.reply_text(text, parse_mode=ParseMode.HTML)
            return
        else:
            raise

    # token usage
    db.set_user_attribute(user_id, "n_generated_images", config.return_n_generated_images + db.get_user_attribute(user_id, "n_generated_images"))

    for i, image_url in enumerate(image_urls):
        await update.message.chat.send_action(action="upload_photo")
        await update.message.reply_photo(image_url, parse_mode=ParseMode.HTML)


async def new_dialog_handle(update: Update, context: CallbackContext):
    await register_user_if_not_exists(update, context, update.message.from_user)
    if await is_previous_message_not_answered_yet(update, context): return

    user_id = update.message.from_user.id
    db.set_user_attribute(user_id, "last_interaction", datetime.now())

    db.start_new_dialog(user_id)
    await update.message.reply_text("Starting new dialog ‚úÖ")

    chat_mode = db.get_user_attribute(user_id, "current_chat_mode")
    await update.message.reply_text(f"{config.chat_modes[chat_mode]['welcome_message']}", parse_mode=ParseMode.HTML)


async def cancel_handle(update: Update, context: CallbackContext):
    await register_user_if_not_exists(update, context, update.message.from_user)

    user_id = update.message.from_user.id
    db.set_user_attribute(user_id, "last_interaction", datetime.now())

    if user_id in user_tasks:
        task = user_tasks[user_id]
        task.cancel()
    else:
        await update.message.reply_text("<i>Nothing to cancel...</i>", parse_mode=ParseMode.HTML)


def is_allow_use_api(user_id: int):
    arr_unlimit = config.unlimit_user
    if arr_unlimit.count(user_id) > 0:
        return True
    daily_stats = db.get_user_daily_stats(user_id, str(date.today()))
    if daily_stats["total_token"] >= config.limit_token or daily_stats["total_chat"] >= config.limit_question:
        return False
    return True
    

def get_chat_mode_menu(page_index: int):
    n_chat_modes_per_page = config.n_chat_modes_per_page
    enabled_chat_modes = {k: v for k, v in config.chat_modes.items() if v.get("enable")}

    text = f"Select <b>chat mode</b> ({len(enabled_chat_modes)} modes available):"

    # buttons
    chat_mode_keys = list(enabled_chat_modes.keys())
    page_chat_mode_keys = chat_mode_keys[page_index * n_chat_modes_per_page:(page_index + 1) * n_chat_modes_per_page]

    keyboard = []
    for chat_mode_key in page_chat_mode_keys:
        name = enabled_chat_modes[chat_mode_key]["name"]
        keyboard.append([InlineKeyboardButton(name, callback_data=f"set_chat_mode|{chat_mode_key}")])

    # pagination
    if len(chat_mode_keys) > n_chat_modes_per_page:
        is_first_page = (page_index == 0)
        is_last_page = ((page_index + 1) * n_chat_modes_per_page >= len(chat_mode_keys))

        if is_first_page:
            keyboard.append([
                InlineKeyboardButton("¬ª", callback_data=f"show_chat_modes|{page_index + 1}")
            ])
        elif is_last_page:
            keyboard.append([
                InlineKeyboardButton("¬´", callback_data=f"show_chat_modes|{page_index - 1}"),
            ])
        else:
            keyboard.append([
                InlineKeyboardButton("¬´", callback_data=f"show_chat_modes|{page_index - 1}"),
                InlineKeyboardButton("¬ª", callback_data=f"show_chat_modes|{page_index + 1}")
            ])

    reply_markup = InlineKeyboardMarkup(keyboard)

    return text, reply_markup


async def show_chat_modes_handle(update: Update, context: CallbackContext):
    await register_user_if_not_exists(update, context, update.message.from_user)
    if await is_previous_message_not_answered_yet(update, context): return

    user_id = update.message.from_user.id
    db.set_user_attribute(user_id, "last_interaction", datetime.now())

    text, reply_markup = get_chat_mode_menu(0)
    await update.message.reply_text(text, reply_markup=reply_markup, parse_mode=ParseMode.HTML)


async def show_chat_modes_callback_handle(update: Update, context: CallbackContext):
     await register_user_if_not_exists(update.callback_query, context, update.callback_query.from_user)
     if await is_previous_message_not_answered_yet(update.callback_query, context): return

     user_id = update.callback_query.from_user.id
     db.set_user_attribute(user_id, "last_interaction", datetime.now())

     query = update.callback_query
     await query.answer()

     page_index = int(query.data.split("|")[1])
     if page_index < 0:
         return

     text, reply_markup = get_chat_mode_menu(page_index)
     try:
         await query.edit_message_text(text, reply_markup=reply_markup, parse_mode=ParseMode.HTML)
     except telegram.error.BadRequest as e:
         if str(e).startswith("Message is not modified"):
             pass


async def set_chat_mode_handle(update: Update, context: CallbackContext):
    await register_user_if_not_exists(update.callback_query, context, update.callback_query.from_user)
    user_id = update.callback_query.from_user.id

    query = update.callback_query
    await query.answer()

    chat_mode = query.data.split("|")[1]

    db.set_user_attribute(user_id, "current_chat_mode", chat_mode)
    db.get_dialog_by_chat_mode(user_id, chat_mode)

    await context.bot.send_message(
        update.callback_query.message.chat.id,
        f"{config.chat_modes[chat_mode]['welcome_message']}",
        parse_mode=ParseMode.HTML
    )


def get_settings_menu(user_id: int):
    current_model = db.get_user_attribute(user_id, "current_model")
    text = config.models["info"][current_model]["description"]

    text += "\n\n"
    score_dict = config.models["info"][current_model]["scores"]
    for score_key, score_value in score_dict.items():
        text += "üü¢" * score_value + "‚ö™Ô∏è" * (5 - score_value) + f" ‚Äì {score_key}\n\n"

    text += "\nSelect <b>model</b>:"

    # buttons to choose models
    buttons = []
    for model_key in config.models["available_text_models"]:
        title = config.models["info"][model_key]["name"]
        if model_key == current_model:
            title = "‚úÖ " + title

        buttons.append(
            InlineKeyboardButton(title, callback_data=f"set_settings|{model_key}")
        )
    reply_markup = InlineKeyboardMarkup([buttons])

    return text, reply_markup


async def settings_handle(update: Update, context: CallbackContext):
    await register_user_if_not_exists(update, context, update.message.from_user)
    if await is_previous_message_not_answered_yet(update, context): return

    user_id = update.message.from_user.id
    db.set_user_attribute(user_id, "last_interaction", datetime.now())

    text, reply_markup = get_settings_menu(user_id)
    await update.message.reply_text(text, reply_markup=reply_markup, parse_mode=ParseMode.HTML)


async def set_settings_handle(update: Update, context: CallbackContext):
    await register_user_if_not_exists(update.callback_query, context, update.callback_query.from_user)
    user_id = update.callback_query.from_user.id

    query = update.callback_query
    await query.answer()

    _, model_key = query.data.split("|")
    db.set_user_attribute(user_id, "current_model", model_key)
    db.start_new_dialog(user_id)

    text, reply_markup = get_settings_menu(user_id)
    try:
        await query.edit_message_text(text, reply_markup=reply_markup, parse_mode=ParseMode.HTML)
    except telegram.error.BadRequest as e:
        if str(e).startswith("Message is not modified"):
            pass


async def show_balance_handle(update: Update, context: CallbackContext):
    await register_user_if_not_exists(update, context, update.message.from_user)
    user_id = update.message.from_user.id
    
    text, reply_markup = payment.get_balance_menu(user_id)
    
    await update.message.reply_text(text, reply_markup=reply_markup, parse_mode=ParseMode.HTML)
    
async def get_contracts(update: Update, context: CallbackContext):
    query = update.callback_query
    await query.answer()
    
    await register_user_if_not_exists(query, context, query.from_user)

    user_id = query.from_user.id
    db.set_user_attribute(user_id, "last_interaction", datetime.now())

    text, reply_markup = payment.get_contracts_menu()
    
    await query.message.edit_text(text,reply_markup=reply_markup, parse_mode=ParseMode.HTML)
    
async def get_providers(update: Update, context: CallbackContext):
    query = update.callback_query
    await query.answer()
    key = query.data.split("|")[1]
    await register_user_if_not_exists(query, context, query.from_user)

    user_id = query.from_user.id
    db.set_user_attribute(user_id, "last_interaction", datetime.now())

    text, reply_markup = payment.get_providers_menu(key)
    
    await query.message.edit_text(text,reply_markup=reply_markup, parse_mode=ParseMode.HTML)

async def send_invoice(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    query = update.callback_query
    await query.answer()
    await register_user_if_not_exists(query, context, query.from_user)
    
    params = query.data.split("|")[1]
    list_params = ast.literal_eval(params)
    payment = config.payments[list_params[0]]
    contract = payment["providers"][list_params[1]]
    
    user_id = query.from_user.id
    db.set_user_attribute(user_id, "last_interaction", datetime.now())
    """Sends an invoice without shipping-payment."""
    chat_id = query.message.chat_id
    title = payment["text"]
    description = payment["description"]
    
    # In order to get a provider_token see https://core.telegram.org/bots/payments#getting-a-token
    currency = "USD"
    # price in dollars
    price = contract["price"]
    # price * 100 so as to include 2 decimal points
    prices = [LabeledPrice(contract["title"], price)]
    id = db.add_new_charge_pending(user_id, chat_id, price, list_params[1], "visa")
    # select a payload just for you to recognize its the donation from your bot
    payload = id
    payment_token = contract["payment_token"]
    # optionally pass need_name=True, need_phone_number=True,
    # need_email=True, need_shipping_address=True, is_flexible=True
    await context.bot.send_invoice(
        chat_id, title, description, payload, payment_token, currency, prices
    )

# after (optional) shipping, it's the pre-checkout
async def precheckout_callback(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Answers the PreQecheckoutQuery"""
    query = update.pre_checkout_query
    # check the payload, is this from your bot?
    #if query.invoice_payload != "Custom-Payload":
    #    # answer False pre_checkout_query
    #    await query.answer(ok=False, error_message="Something went wrong...")
    #else:
    id = query.invoice_payload
    if db.check_if_invoice_available(id):
        await query.answer(ok=True)
    else:
        await query.answer(ok=False)


# finally, after contacting the payment provider...
async def successful_payment_callback(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Confirms the successful payment."""
    # do something after successfully receiving payment?
    print("payload: " + update.message.successful_payment.invoice_payload)
    db.update_charge_success(update.message.successful_payment.invoice_payload)
    await update.message.reply_text("Thank you for your payment!")

async def edited_message_handle(update: Update, context: CallbackContext):
    if update.edited_message.chat.type == "private":
        text = "ü•≤ Unfortunately, message <b>editing</b> is not supported"
        await update.edited_message.reply_text(text, parse_mode=ParseMode.HTML)


async def error_handle(update: Update, context: CallbackContext) -> None:
    logger.error(msg="Exception while handling an update:", exc_info=context.error)

    try:
        # collect error message
        tb_list = traceback.format_exception(None, context.error, context.error.__traceback__)
        tb_string = "".join(tb_list)
        update_str = update.to_dict() if isinstance(update, Update) else str(update)
        message = (
            f"An exception was raised while handling an update\n"
            f"<pre>update = {html.escape(json.dumps(update_str, indent=2, ensure_ascii=False))}"
            "</pre>\n\n"
            f"<pre>{html.escape(tb_string)}</pre>"
        )

        # split text into multiple messages due to 4096 character limit
        for message_chunk in split_text_into_chunks(message, 4096):
            try:
                await context.bot.send_message(update.effective_chat.id, message_chunk, parse_mode=ParseMode.HTML)
            except telegram.error.BadRequest:
                # answer has invalid characters, so we send it without parse_mode
                await context.bot.send_message(update.effective_chat.id, message_chunk)
    except:
        await context.bot.send_message(update.effective_chat.id, "Some error in error handler")

async def post_init(application: Application):
    await application.bot.set_my_commands([
        BotCommand("/new", "Start new dialog"),
        BotCommand("/mode", "Select chat mode"),
        BotCommand("/retry", "Re-generate response for previous query"),
        BotCommand("/balance", "Show balance"),
        BotCommand("/settings", "Show settings"),
        BotCommand("/help", "Show help message"),
    ])

def run_bot() -> None:
    application = (
        ApplicationBuilder()
        .token(config.telegram_token)
        .concurrent_updates(True)
        .rate_limiter(AIORateLimiter(max_retries=5))
        .http_version("1.1")
        .get_updates_http_version("1.1")
        .post_init(post_init)
        .build()
    )

    # add handlers
    user_filter = filters.ALL
    if len(config.allowed_telegram_usernames) > 0:
        usernames = [x for x in config.allowed_telegram_usernames if isinstance(x, str)]
        any_ids = [x for x in config.allowed_telegram_usernames if isinstance(x, int)]
        user_ids = [x for x in any_ids if x > 0]
        group_ids = [x for x in any_ids if x < 0]
        user_filter = filters.User(username=usernames) | filters.User(user_id=user_ids) | filters.Chat(chat_id=group_ids)

    application.add_handler(CommandHandler("start", start_handle, filters=user_filter))
    application.add_handler(CommandHandler("help", help_handle, filters=user_filter))
    application.add_handler(CommandHandler("help_group_chat", help_group_chat_handle, filters=user_filter))

    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND & user_filter, message_handle))
    application.add_handler(CommandHandler("retry", retry_handle, filters=user_filter))
    application.add_handler(CommandHandler("new", new_dialog_handle, filters=user_filter))
    application.add_handler(CommandHandler("cancel", cancel_handle, filters=user_filter))

    application.add_handler(MessageHandler(filters.VOICE & user_filter, voice_message_handle))
    application.add_handler(MessageHandler(filters.Document.IMAGE & user_filter, image_message_handle))
    application.add_handler(MessageHandler(filters.Document.TEXT & user_filter, text_file_message_handle))
    application.add_handler(MessageHandler(filters.Document.VIDEO & user_filter, video_message_handle))
    # application.add_handler(MessageHandler(filters.Document.MIME_TYPE == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' & user_filter, ppt_message_handle))
    application.add_handler(MessageHandler(filters.Document.DOCX & user_filter, doc_message_handle))
    application.add_handler(MessageHandler(filters.Document.PDF & user_filter, pdf_message_handle))

    application.add_handler(CommandHandler("mode", show_chat_modes_handle, filters=user_filter))
    application.add_handler(CallbackQueryHandler(show_chat_modes_callback_handle, pattern="^show_chat_modes"))
    application.add_handler(CallbackQueryHandler(set_chat_mode_handle, pattern="^set_chat_mode"))

    application.add_handler(CommandHandler("settings", settings_handle, filters=user_filter))
    application.add_handler(CallbackQueryHandler(set_settings_handle, pattern="^set_settings"))

    application.add_handler(CommandHandler("balance", show_balance_handle, filters=user_filter))
    
    #subscription
    application.add_handler(CallbackQueryHandler(get_contracts, pattern="^get_contracts"))
    application.add_handler(CallbackQueryHandler(get_providers, pattern="^get_providers"))
    application.add_handler(CallbackQueryHandler(send_invoice, pattern="^send_invoice"))
    application.add_handler(CallbackQueryHandler(payment.back_balance_menu, pattern="^back_balance_menu"))
    application.add_handler(CallbackQueryHandler(payment.back_contracts_menu, pattern="^back_contracts_menu"))
    
    # Pre-checkout handler to final check
    application.add_handler(PreCheckoutQueryHandler(precheckout_callback))

    # Success! Notify your user!
    application.add_handler(
        MessageHandler(filters.SUCCESSFUL_PAYMENT, successful_payment_callback)
    )

    application.add_error_handler(error_handle)

    # start the bot
    application.run_polling()


if __name__ == "__main__":
    run_bot()
